#!/usr/bin/env python
"""
Generate Comprehensive MyQA Agent PowerPoint Presentation
Includes functional flow diagrams, API documentation, and development details
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_comprehensive_presentation():
    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Define colors
    PRIMARY_COLOR = RGBColor(13, 110, 253)  # Blue
    SECONDARY_COLOR = RGBColor(25, 135, 84)  # Green
    ACCENT_COLOR = RGBColor(220, 53, 69)  # Red
    WARNING_COLOR = RGBColor(255, 193, 7)  # Yellow
    TEXT_COLOR = RGBColor(33, 37, 41)  # Dark gray
    LIGHT_BG = RGBColor(248, 249, 250)  # Light gray
    
    def add_title_slide(title, subtitle=""):
        """Add title slide"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = PRIMARY_COLOR
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
        title_frame = title_box.text_frame
        title_frame.text = title
        title_frame.paragraphs[0].font.size = Pt(54)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Subtitle
        if subtitle:
            subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(2))
            subtitle_frame = subtitle_box.text_frame
            subtitle_frame.word_wrap = True
            subtitle_frame.text = subtitle
            subtitle_frame.paragraphs[0].font.size = Pt(28)
            subtitle_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
            subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    def add_content_slide(title, content_points, bg_color=RGBColor(255, 255, 255)):
        """Add content slide with bullet points"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = bg_color
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
        title_frame = title_box.text_frame
        title_frame.text = title
        title_frame.paragraphs[0].font.size = Pt(40)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = PRIMARY_COLOR
        
        # Separator line
        line = slide.shapes.add_shape(1, Inches(0.5), Inches(1.3), Inches(9), Inches(0))
        line.line.color.rgb = SECONDARY_COLOR
        line.line.width = Pt(3)
        
        # Content
        content_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.6), Inches(8.6), Inches(5.4))
        text_frame = content_box.text_frame
        text_frame.word_wrap = True
        
        for i, point in enumerate(content_points):
            if i > 0:
                text_frame.add_paragraph()
            p = text_frame.paragraphs[i]
            p.text = point
            p.font.size = Pt(16)
            p.font.color.rgb = TEXT_COLOR
            p.space_before = Pt(6)
            p.space_after = Pt(6)
            p.level = 0
    
    def add_flowchart_slide(title):
        """Add flowchart slide"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(255, 255, 255)
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.7))
        title_frame = title_box.text_frame
        title_frame.text = title
        title_frame.paragraphs[0].font.size = Pt(38)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = PRIMARY_COLOR
        
        # Separator line
        line = slide.shapes.add_shape(1, Inches(0.5), Inches(1.2), Inches(9), Inches(0))
        line.line.color.rgb = SECONDARY_COLOR
        line.line.width = Pt(3)
        
        return slide
    
    def add_box_with_text(slide, left, top, width, height, text, bg_color, text_color=RGBColor(255, 255, 255)):
        """Add a box with text"""
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, 
                                      Inches(left), Inches(top), Inches(width), Inches(height))
        shape.fill.solid()
        shape.fill.fore_color.rgb = bg_color
        shape.line.color.rgb = RGBColor(100, 100, 100)
        shape.line.width = Pt(2)
        
        text_frame = shape.text_frame
        text_frame.clear()
        text_frame.word_wrap = True
        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = text_frame.paragraphs[0]
        p.text = text
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = text_color
        p.alignment = PP_ALIGN.CENTER
        
        return shape
    
    def add_arrow(slide, x1, y1, x2, y2):
        """Add arrow connector"""
        arrow = slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
        arrow.line.color.rgb = SECONDARY_COLOR
        arrow.line.width = Pt(2)
    
    # ===== PAGE 1: Title Slide =====
    add_title_slide(
        "MyQA Agent",
        "AI-Powered Automated Test Case Generator\nComprehensive Documentation & Architecture"
    )
    
    # ===== PAGE 2: Project Overview =====
    add_content_slide(
        "📋 Project Overview",
        [
            "Project Name: MyQA Agent (Agentic QA Application)",
            "",
            "Purpose:",
            "  • Automate test case generation from business requirements",
            "  • Convert documents (PDF, DOCX, TXT) into comprehensive test cases",
            "  • Support positive, negative, and boundary test scenarios",
            "",
            "Key Capabilities:",
            "  ✓ Document Upload & Parsing (Multi-format support)",
            "  ✓ AI-Powered Requirement Analysis (Extract user stories, criteria)",
            "  ✓ Intelligent Test Case Generation (Auto-generate test steps)",
            "  ✓ Multi-Format Export (Word .docx, Excel .xlsx)",
            "",
            "Technology Stack:",
            "  Backend: FastAPI, Python, Pydantic  |  Frontend: HTML5, CSS3, JavaScript",
            "  Processing: PyPDF2, python-docx  |  Export: openpyxl, python-docx"
        ]
    )
    
    # ===== PAGE 3: System Architecture =====
    add_content_slide(
        "🏗️ System Architecture",
        [
            "Multi-Layer Architecture:",
            "",
            "1️⃣  Presentation Layer (Frontend)",
            "   • HTML5 UI with responsive design (Bootstrap 5.3)",
            "   • JavaScript for interactive features & drag-drop upload",
            "   • Real-time progress feedback & status updates",
            "",
            "2️⃣  API Layer (FastAPI Routers)",
            "   • /api/upload/document - File upload & validation",
            "   • /api/analyze/document - Content analysis & extraction",
            "   • /api/generate/test-cases - Intelligent test generation",
            "   • /api/export/test-cases - Multi-format export",
            "",
            "3️⃣  Business Logic Layer (Services)",
            "   • DocumentParser: Extract text from various formats",
            "   • AIAgent: Generate test cases using intelligent patterns",
            "",
            "4️⃣  Data Layer (File System)",
            "   • /uploads/ - Stores uploaded documents",
            "   • /exports/ - Stores generated test case files"
        ]
    )
    
    # ===== PAGE 4: Functional Flow Diagram =====
    slide = add_flowchart_slide("🔄 End-to-End Functional Flow")
    
    # Flow boxes
    add_box_with_text(slide, 1.5, 1.5, 2.5, 0.8, "1. User Uploads\nDocument", PRIMARY_COLOR)
    add_arrow(slide, 4, 1.9, 4.7, 1.9)
    add_box_with_text(slide, 4.7, 1.5, 2.5, 0.8, "2. Upload API\nValidation", SECONDARY_COLOR)
    add_arrow(slide, 7.2, 1.9, 7.9, 1.9)
    add_box_with_text(slide, 7.9, 1.5, 1.5, 0.8, "Save to\n/uploads/", WARNING_COLOR, TEXT_COLOR)
    
    add_arrow(slide, 3.2, 2.3, 3.2, 2.8)
    add_box_with_text(slide, 1.5, 2.8, 3.4, 0.8, "3. Analyze API\nParse & Extract Content", PRIMARY_COLOR)
    add_arrow(slide, 4.9, 3.2, 5.5, 3.2)
    add_box_with_text(slide, 5.5, 2.8, 3.4, 0.8, "4. AI Processing\nIdentify Requirements", SECONDARY_COLOR)
    
    add_arrow(slide, 3.2, 3.6, 3.2, 4.1)
    add_box_with_text(slide, 1.5, 4.1, 3.4, 0.8, "5. Generate API\nCreate Test Cases", PRIMARY_COLOR)
    add_arrow(slide, 4.9, 4.5, 5.5, 4.5)
    add_box_with_text(slide, 5.5, 4.1, 3.4, 0.8, "6. Export API\nFormat Output (DOCX/XLSX)", SECONDARY_COLOR)
    
    add_arrow(slide, 3.2, 4.9, 3.2, 5.4)
    add_box_with_text(slide, 1.5, 5.4, 6.9, 0.8, "7. Download\nUser Receives Generated Test Cases", WARNING_COLOR, TEXT_COLOR)
    
    # ===== PAGE 5: API Details - Upload =====
    add_content_slide(
        "📤 API #1: Upload Document Endpoint",
        [
            "Endpoint: POST /api/upload/document",
            "",
            "Purpose: Accept and validate document uploads",
            "",
            "Request Parameters:",
            "  • file (FormData): PDF, DOCX, or TXT file (Max 50MB)",
            "",
            "Response Structure:",
            "  {",
            "    \"success\": true,",
            "    \"file_id\": \"20260407_235207_sample_requirements.txt\",",
            "    \"file_name\": \"sample_requirements.txt\",",
            "    \"file_size\": 4783,",
            "    \"upload_time\": \"2026-04-07T23:52:07\",",
            "    \"message\": \"File uploaded successfully\"",
            "  }",
            "",
            "Validation Logic:",
            "  ✓ File type check (.pdf, .docx, .txt)",
            "  ✓ File size validation (50MB limit)",
            "  ✓ Unique file ID generation with timestamp"
        ]
    )
    
    # ===== PAGE 6: API Details - Analyze =====
    add_content_slide(
        "🔍 API #2: Analyze Document Endpoint",
        [
            "Endpoint: POST /api/analyze/document",
            "",
            "Purpose: Extract structured information from documents",
            "",
            "Request Parameters:",
            "  • file_id (string): ID from upload response",
            "",
            "Response Structure:",
            "  {",
            "    \"success\": true,",
            "    \"file_id\": \"...\",",
            "    \"analysis\": {",
            "      \"user_stories\": [\"As a user...\", ...],",
            "      \"acceptance_criteria\": [\"Given...\", ...],",
            "      \"business_rules\": [\"Email must be unique...\", ...],",
            "      \"constraints\": [\"Must complete in 3 seconds...\", ...],",
            "      \"preconditions\": [\"User logged in...\", ...]",
            "    }",
            "  }",
            "",
            "Processing Steps:",
            "  1. Load file from /uploads/ directory",
            "  2. Extract raw text based on file type",
            "  3. Parse and identify requirement patterns",
            "  4. Categorize into user stories, criteria, rules, constraints"
        ]
    )
    
    # ===== PAGE 7: API Details - Generate =====
    add_content_slide(
        "⚡ API #3: Generate Test Cases Endpoint",
        [
            "Endpoint: POST /api/generate/test-cases",
            "",
            "Purpose: Intelligently generate test cases from requirements",
            "",
            "Request Parameters:",
            "  • file_id (string): File ID from upload",
            "  • num_test_cases (int): 1-20 test cases (default: 5)",
            "  • include_positive (bool): Include positive scenarios (default: true)",
            "  • include_negative (bool): Include negative scenarios (default: true)",
            "  • include_boundary (bool): Include boundary scenarios (default: true)",
            "",
            "Response Structure:",
            "  {",
            "    \"success\": true,",
            "    \"total_count\": 10,",
            "    \"test_cases\": [",
            "      {",
            "        \"test_case_id\": \"TC_POS_001\",",
            "        \"name\": \"Positive Test Case\",",
            "        \"test_type\": \"positive\",",
            "        \"priority\": \"high\",",
            "        \"test_steps\": [{\"step_number\": 1, \"action\": \"...\"}],",
            "        \"expected_result\": \"...\"",
            "      }",
            "    ]",
            "  }"
        ]
    )
    
    # ===== PAGE 8: API Details - Export =====
    add_content_slide(
        "💾 API #4: Export Test Cases Endpoint",
        [
            "Endpoint: POST /api/export/test-cases",
            "",
            "Purpose: Export test cases to Word or Excel format",
            "",
            "Request Parameters:",
            "  • file_id (string): File ID from upload",
            "  • format (string): \"docx\" or \"xlsx\"",
            "  • num_test_cases (int): Number of test cases to export",
            "",
            "Response Structure:",
            "  • Returns binary file (application/vnd.openxmlformats-officedocument.wordprocessingml.document)",
            "  • Returns binary file (application/vnd.openxmlformats-officedocument.spreadsheetml.sheet)",
            "",
            "Export Details:",
            "",
            "📄 Word (.docx) Export:",
            "  • Formatted document with headers, sections, tables",
            "  • Includes test case details, steps, and results",
            "  • Professional formatting with colors and styling",
            "",
            "📊 Excel (.xlsx) Export:",
            "  • Spreadsheet with columns: ID, Name, Type, Priority, Steps, Results",
            "  • Each test case on separate row",
            "  • Auto-sized columns and conditional formatting"
        ]
    )
    
    # ===== PAGE 9: Development Flow - Backend Setup =====
    add_content_slide(
        "🔧 Development Flow - Backend Setup",
        [
            "Phase 1: Project Initialization",
            "",
            "1. Environment Setup:",
            "   • Create Python virtual environment (.venv)",
            "   • Install dependencies: FastAPI, Uvicorn, PyPDF2, python-docx, openpyxl",
            "",
            "2. Project Structure:",
            "   app/",
            "   ├── main.py (FastAPI application)",
            "   ├── api/ (Route handlers)",
            "   │  ├── upload.py",
            "   │  ├── analyze.py",
            "   │  ├── generate.py",
            "   │  └── export.py",
            "   ├── services/ (Business logic)",
            "   │  ├── document_parser.py",
            "   │  └── ai_agent.py",
            "   ├── models/ (Data models)",
            "   │  └── testcase.py",
            "   └── utils/ (Utilities)",
            "      ├── excel_exporter.py",
            "      └── word_exporter.py"
        ]
    )
    
    # ===== PAGE 10: Development Flow - Core Services =====
    add_content_slide(
        "🧠 Development Flow - Core Services",
        [
            "Phase 2: Core Service Implementation",
            "",
            "1. DocumentParser Service (document_parser.py):",
            "   • Reads PDF files using PyPDF2",
            "   • Extracts DOCX content using python-docx",
            "   • Reads plain text files directly",
            "   • Returns unified text format for processing",
            "",
            "2. AIAgent Service (ai_agent.py):",
            "   • Analyzes extracted text using pattern recognition",
            "   • Identifies user stories, acceptance criteria",
            "   • Generates positive, negative, boundary test cases",
            "   • Creates detailed test steps with expected results",
            "",
            "3. TestCase Model (models/testcase.py):",
            "   • Pydantic data validation models",
            "   • TestStep, TestCase, AnalysisResult schemas",
            "   • Type hints and field validation",
            "",
            "4. Exporters (excel_exporter.py, word_exporter.py):",
            "   • Format test cases for Word documents",
            "   • Format test cases for Excel spreadsheets",
            "   • Apply professional styling and formatting"
        ]
    )
    
    # ===== PAGE 11: Development Flow - API Implementation =====
    add_content_slide(
        "🛣️ Development Flow - API Implementation",
        [
            "Phase 3: API Endpoint Development",
            "",
            "1. Upload Endpoint (app/api/upload.py):",
            "   • FastAPI router with POST method",
            "   • FormData file parameter",
            "   • File validation and storage logic",
            "   • Returns file_id for tracking",
            "",
            "2. Analyze Endpoint (app/api/analyze.py):",
            "   • Receives file_id from upload",
            "   • Calls DocumentParser service",
            "   • Processes content with AIAgent",
            "   • Returns structured analysis data",
            "",
            "3. Generate Endpoint (app/api/generate.py):",
            "   • Receives file_id and generation parameters",
            "   • Calls AIAgent to generate test cases",
            "   • Filters by positive/negative/boundary flags",
            "   • Returns list of test cases",
            "",
            "4. Export Endpoint (app/api/export.py):",
            "   • Receives file_id and format parameter",
            "   • Generates test cases if not already done",
            "   • Formats output (Word or Excel)",
            "   • Returns binary file for download"
        ]
    )
    
    # ===== PAGE 12: Development Flow - Frontend =====
    add_content_slide(
        "🎨 Development Flow - Frontend Development",
        [
            "Phase 4: User Interface Implementation",
            "",
            "1. HTML Structure (public/index.html):",
            "   • Bootstrap 5.3 responsive layout",
            "   • Upload area with drag-drop support",
            "   • Analysis results display",
            "   • Generation settings (# test cases, types)",
            "   • Test cases list with detail view",
            "   • Export buttons (Word, Excel)",
            "",
            "2. Styling (public/style.css):",
            "   • Custom CSS variables for colors",
            "   • Card-based design pattern",
            "   • Responsive grid layout (col-lg-4, col-lg-8)",
            "   • Progress bars and loading spinners",
            "   • Professional color scheme",
            "",
            "3. JavaScript Logic (public/script.js):",
            "   • Event listeners for all UI interactions",
            "   • Fetch API calls to backend endpoints",
            "   • File validation on client side",
            "   • Progress tracking and UI updates",
            "   • Error handling and user feedback"
        ]
    )
    
    # ===== PAGE 13: Development Flow - Integration & Testing =====
    add_content_slide(
        "✅ Development Flow - Integration & Testing",
        [
            "Phase 5: Integration & Quality Assurance",
            "",
            "1. CORS Configuration:",
            "   • Enable cross-origin requests from frontend to API",
            "   • Configure allowed methods: GET, POST, OPTIONS",
            "   • Allow credentials and custom headers",
            "",
            "2. Error Handling:",
            "   • 400: Invalid file format or missing parameters",
            "   • 404: File not found in /uploads/ directory",
            "   • 500: Server-side errors with detailed messages",
            "   • Custom exception handlers in FastAPI",
            "",
            "3. Testing Workflow:",
            "   ✓ Upload various file formats (PDF, DOCX, TXT)",
            "   ✓ Analyze extracted content accuracy",
            "   ✓ Generate test cases with different options",
            "   ✓ Export to Word and Excel formats",
            "   ✓ Verify downloaded files are valid",
            "",
            "4. Performance Optimization:",
            "   • File caching to avoid re-parsing",
            "   • Efficient text processing algorithms",
            "   • Streaming responses for large files"
        ]
    )
    
    # ===== PAGE 14: Deployment Architecture =====
    add_content_slide(
        "🚀 Deployment & Scalability",
        [
            "Deployment Options:",
            "",
            "1. Local Development:",
            "   uvicorn app.main:app --reload --host 0.0.0.0 --port 8000",
            "",
            "2. Production Server Options:",
            "   • Gunicorn with Uvicorn workers: gunicorn -w 4 -k uvicorn.workers.UvicornWorker app.main:app",
            "   • Docker containerization for consistency",
            "   • Environment variables for configuration",
            "",
            "3. Scalability Features:",
            "   ✓ Async/await for concurrent request handling",
            "   ✓ Stateless API design for horizontal scaling",
            "   ✓ File-based storage (upgradeable to cloud storage)",
            "   ✓ Load balancing support",
            "",
            "4. File Organization:",
            "   /uploads/ → Stores user-uploaded documents",
            "   /exports/ → Stores generated test case files",
            "   .venv/ → Python virtual environment",
            "   /public/ → Frontend assets (HTML, CSS, JS)"
        ]
    )
    
    # ===== PAGE 15: Key Technologies & Dependencies =====
    add_content_slide(
        "📦 Technologies & Dependencies",
        [
            "Backend Framework:",
            "  • FastAPI (0.104.1): Modern async web framework with auto-docs",
            "  • Uvicorn (0.24.0): ASGI server for running FastAPI",
            "  • Pydantic (2.5.0): Data validation using Python type hints",
            "",
            "Document Processing:",
            "  • PyPDF2 (3.0.1): PDF text extraction",
            "  • python-docx (0.8.11): DOCX file parsing and creation",
            "  • openpyxl (3.1.2): Excel file creation and formatting",
            "",
            "Frontend:",
            "  • HTML5: Semantic markup and structure",
            "  • CSS3: Modern styling with variables and flexbox",
            "  • Bootstrap 5.3.0: Responsive design framework",
            "  • JavaScript ES6+: Interactive features and API calls",
            "",
            "Presentation:",
            "  • python-pptx: PowerPoint generation programmatically"
        ]
    )
    
    # ===== PAGE 16: Code Quality & Best Practices =====
    add_content_slide(
        "⭐ Code Quality & Best Practices",
        [
            "1. Architecture Principles:",
            "   ✓ Separation of concerns (API, Services, Models, Utils)",
            "   ✓ Modular router design for maintainability",
            "   ✓ Reusable service classes",
            "   ✓ Type hints throughout codebase",
            "",
            "2. Error Handling:",
            "   ✓ Try-catch blocks in critical operations",
            "   ✓ Meaningful error messages for debugging",
            "   ✓ Custom exception handlers",
            "",
            "3. Documentation:",
            "   ✓ Docstrings for all functions and classes",
            "   ✓ Inline comments for complex logic",
            "   ✓ API documentation via FastAPI /docs",
            "",
            "4. Code Organization:",
            "   ✓ Consistent naming conventions",
            "   ✓ DRY principle (Don't Repeat Yourself)",
            "   ✓ Configuration management via environment variables",
            "   ✓ Logging for debugging and monitoring"
        ]
    )
    
    # ===== PAGE 17: Future Enhancements =====
    add_content_slide(
        "🎯 Future Enhancements & Roadmap",
        [
            "Short Term (Next 1-3 months):",
            "  • Database integration (PostgreSQL/MongoDB)",
            "  • User authentication and authorization",
            "  • API rate limiting and throttling",
            "  • Advanced file format support (XML, JSON)",
            "",
            "Medium Term (3-6 months):",
            "  • Machine learning model for smarter test generation",
            "  • Test case templates library",
            "  • Bulk file upload and processing",
            "  • Collaboration features (multiple users)",
            "",
            "Long Term (6+ months):",
            "  • Mobile app for iOS/Android",
            "  • Integration with test management tools (Jira, TestRail)",
            "  • CI/CD pipeline integration",
            "  • Real-time collaboration like Google Docs",
            "  • Advanced analytics and reporting"
        ]
    )
    
    # ===== PAGE 18: Summary & Key Takeaways =====
    add_content_slide(
        "📌 Summary & Key Takeaways",
        [
            "What is MyQA Agent?",
            "  An intelligent AI-powered tool that converts business requirements",
            "  into comprehensive test cases automatically.",
            "",
            "Key Features:",
            "  ✓ Multi-format document upload (PDF, DOCX, TXT)",
            "  ✓ Intelligent requirement extraction and analysis",
            "  ✓ Smart test case generation (positive/negative/boundary)",
            "  ✓ Professional export formats (Word, Excel)",
            "  ✓ User-friendly web interface with drag-drop upload",
            "",
            "Architecture Highlights:",
            "  ✓ Modern FastAPI backend with async support",
            "  ✓ Responsive HTML5/CSS3/JavaScript frontend",
            "  ✓ Modular service-oriented design",
            "  ✓ Scalable and maintainable codebase",
            "",
            "Value Proposition:",
            "  ⏱️  Save hours of manual test case writing",
            "  💰 Reduce QA effort and associated costs",
            "  📊 Ensure comprehensive test coverage",
            "  🎯 Improve testing consistency and quality"
        ]
    )
    
    # Save presentation
    output_file = "MyQA_Agent_Comprehensive_Presentation.pptx"
    prs.save(output_file)
    print(f"✅ Comprehensive presentation created: {output_file}")
    print(f"📊 Pages: {len(prs.slides)}")
    print(f"📁 Location: {output_file}")
    print(f"\nPresentation includes:")
    print(f"  • System architecture & overview")
    print(f"  • Functional flow diagrams")
    print(f"  • Complete API documentation (4 endpoints)")
    print(f"  • Development flow breakdown (5 phases)")
    print(f"  • Technology stack & dependencies")
    print(f"  • Best practices & code quality")
    print(f"  • Future roadmap & enhancements")

if __name__ == "__main__":
    create_comprehensive_presentation()
