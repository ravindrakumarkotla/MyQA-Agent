#!/usr/bin/env python
"""
Generate MyQA Agent Team Demo PowerPoint Presentation (5 Pages)
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_presentation():
    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Define colors
    PRIMARY_COLOR = RGBColor(13, 110, 253)  # Blue
    SECONDARY_COLOR = RGBColor(25, 135, 84)  # Green
    ACCENT_COLOR = RGBColor(220, 53, 69)  # Red
    TEXT_COLOR = RGBColor(33, 37, 41)  # Dark gray
    
    def add_title_slide(title, subtitle=""):
        """Add title slide"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = PRIMARY_COLOR
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
        title_frame = title_box.text_frame
        title_frame.text = title
        title_frame.paragraphs[0].font.size = Pt(54)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Subtitle
        if subtitle:
            subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(2))
            subtitle_frame = subtitle_box.text_frame
            subtitle_frame.word_wrap = True
            subtitle_frame.text = subtitle
            subtitle_frame.paragraphs[0].font.size = Pt(32)
            subtitle_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
            subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    def add_content_slide(title, content_points, bg_color=RGBColor(255, 255, 255)):
        """Add content slide with bullet points"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = bg_color
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
        title_frame = title_box.text_frame
        title_frame.text = title
        title_frame.paragraphs[0].font.size = Pt(40)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = PRIMARY_COLOR
        
        # Separator line
        line = slide.shapes.add_shape(1, Inches(0.5), Inches(1.3), Inches(9), Inches(0))
        line.line.color.rgb = SECONDARY_COLOR
        line.line.width = Pt(3)
        
        # Content
        content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(8.4), Inches(5.4))
        text_frame = content_box.text_frame
        text_frame.word_wrap = True
        
        for i, point in enumerate(content_points):
            if i > 0:
                text_frame.add_paragraph()
            p = text_frame.paragraphs[i]
            p.text = point
            p.font.size = Pt(18)
            p.font.color.rgb = TEXT_COLOR
            p.space_before = Pt(8)
            p.space_after = Pt(8)
            p.level = 0
    
    # ===== PAGE 1: Title Slide =====
    add_title_slide(
        "MyQA Agent",
        "AI-Powered Automated Test Case Generator\nComplete Project Journey"
    )
    
    # ===== PAGE 2: End-to-End Project Journey =====
    add_content_slide(
        "📊 Project Completion Journey",
        [
            "✓ Phase 1: Backend Development",
            "  • FastAPI REST API (4 endpoints)",
            "  • Document Processing: PDF, DOCX, TXT parsing",
            "  • Test Generation Engine (Positive/Negative/Boundary)",
            "  • Database: File-system based storage with auto-indexing",
            "",
            "✓ Phase 2: Frontend Development",
            "  • Professional Web UI (HTML5, CSS3, Bootstrap)",
            "  • Real-time interactive features (Drag & Drop, Live updates)",
            "  • Export to Word & Excel formats",
            "",
            "✓ Phase 3: Documentation & Deployment",
            "  • 16+ comprehensive documentation files",
            "  • Production-ready deployment options",
            "  • Complete API & Architecture documentation"
        ]
    )
    
    # ===== PAGE 3: Application Details & User Interaction =====
    add_content_slide(
        "🎯 MyQA Agent Application Details",
        [
            "What is MyQA Agent?",
            "  Intelligent QA tool that converts business requirements",
            "  into comprehensive test cases automatically",
            "",
            "Key Features:",
            "  🔼 Upload Documents (PDF, DOCX, TXT files)",
            "  🔍 Auto-Extract: User Stories, Acceptance Criteria,",
            "     Business Rules, Constraints",
            "  ✨ Smart Test Generation: Positive, Negative, Boundary tests",
            "  📥 Multi-Format Export: Word (.docx) & Excel (.xlsx)",
            "  🎨 Responsive UI: Works on Desktop, Tablet, Mobile",
            "",
            "Tech Stack:",
            "  Backend: FastAPI, Python, Pydantic",
            "  Frontend: HTML5, CSS3, Bootstrap, JavaScript",
            "  Storage: File-system (uploads/, exports/ directories)"
        ]
    )
    
    # ===== PAGE 4: Technical Flowchart - QA Engineer Journey =====
    add_content_slide(
        "⚙️ End-to-End Technical Flowchart",
        [
            "QA Engineer Workflow:",
            "",
            "1️⃣  UPLOAD PHASE:",
            "   User uploads requirements file → Validation → Save to /uploads/",
            "",
            "2️⃣  ANALYSIS PHASE:",
            "   Parse document → Extract text → NLP analysis →",
            "   Identify requirements → Cache in memory",
            "",
            "3️⃣  GENERATION PHASE:",
            "   Process each requirement → Generate test steps →",
            "   Create positive/negative/boundary scenarios →",
            "   Assign priority & categorize",
            "",
            "4️⃣  EXPORT PHASE:",
            "   Format test cases → Generate Word/Excel →",
            "   Return binary file → Auto-download to user machine",
            "",
            "🔄 Data Flow:",
            "   User → UI → FastAPI → Parser → AI Engine → Exporter → File"
        ]
    )
    
    # ===== PAGE 5: Key Achievements & Value Proposition =====
    add_content_slide(
        "🏆 Key Achievements & Business Value",
        [
            "📈 Project Metrics:",
            "  • 2,000+ lines of production-ready code",
            "  • 10,000+ lines of comprehensive documentation",
            "  • 100+ code examples for developers",
            "  • 26+ troubleshooting scenarios covered",
            "",
            "💼 Business Value:",
            "  ⏱️  Time Saving: Convert hours of manual test writing",
            "      into minutes of automated generation",
            "  💰 Cost Reduction: Minimize QA effort and errors",
            "  📊 Quality Improvement: Comprehensive test coverage",
            "  🔄 Consistency: Standardized test case format",
            "",
            "🚀 Deployment Ready:",
            "  • Docker containerization supported",
            "  • Multiple server options (Gunicorn, Uvicorn)",
            "  • Scalable architecture for team use",
            "  • Production-grade error handling & logging"
        ]
    )
    
    # Save presentation
    output_file = "MyQA_Agent_Team_Demo.pptx"
    prs.save(output_file)
    print(f"✅ Presentation created: {output_file}")
    print(f"📊 Pages: 5")
    print(f"📁 Location: {output_file}")

if __name__ == "__main__":
    create_presentation()
